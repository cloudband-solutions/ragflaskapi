from pathlib import Path
import io
import json
import os
from urllib.parse import unquote_plus

import boto3
import click
from flask import current_app
from flask.cli import with_appcontext
from flask_migrate.cli import db as db_cli
from openai import OpenAI
from pypdf import PdfReader
from sqlalchemy import create_engine, text
from sqlalchemy.engine import make_url
import tiktoken

from app import db
from app.helpers.api_helpers import build_password_hash
from app.models.document import Document
from app.models.document_embedding import DocumentEmbedding
from app.models.user import User
from app.storage import get_storage


DEFAULT_CHUNK_TOKENS = 800
DEFAULT_CHUNK_OVERLAP = 100
DEFAULT_BATCH_SIZE = 50
DEFAULT_SQS_WAIT_TIME = 10


def _validate_chunking_options(chunk_size, chunk_overlap):
    if chunk_size <= 0:
        raise click.ClickException("--chunk-size must be greater than 0.")
    if chunk_overlap < 0:
        raise click.ClickException("--chunk-overlap must be 0 or greater.")
    if chunk_overlap >= chunk_size:
        raise click.ClickException("--chunk-overlap must be less than --chunk-size.")


def _openai_client_and_model():
    use_openai = current_app.config.get("USE_OPENAI", "true")
    if str(use_openai).lower() not in {"1", "true", "yes", "y"}:
        raise click.ClickException("USE_OPENAI is disabled; no OpenAI embedder available.")

    api_key = current_app.config.get("OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY")
    model = current_app.config.get("OPENAI_EMBEDDING_MODEL") or os.getenv(
        "OPENAI_EMBEDDING_MODEL"
    )
    if not api_key:
        raise click.ClickException("OPENAI_API_KEY is required.")
    if not model:
        raise click.ClickException("OPENAI_EMBEDDING_MODEL is required.")

    return OpenAI(api_key=api_key), model


def _local_embedder():
    model_path = current_app.config.get("LOCAL_EMBEDDING_MODEL_PATH") or os.getenv(
        "LOCAL_EMBEDDING_MODEL_PATH"
    )
    if not model_path:
        raise click.ClickException("LOCAL_EMBEDDING_MODEL_PATH is required.")
    if not os.path.exists(model_path):
        raise click.ClickException(f"Local model not found: {model_path}")
    if not model_path.lower().endswith(".gguf"):
        click.echo("Warning: local model does not look like a .gguf file.")

    try:
        from llama_cpp import Llama
    except ImportError as exc:
        raise click.ClickException(
            "llama-cpp-python is not installed. Add it to requirements and install."
        ) from exc

    n_ctx = int(
        current_app.config.get("LOCAL_EMBEDDING_N_CTX", 2048)
        or os.getenv("LOCAL_EMBEDDING_N_CTX", "2048")
    )
    n_threads = int(
        current_app.config.get("LOCAL_EMBEDDING_N_THREADS", 4)
        or os.getenv("LOCAL_EMBEDDING_N_THREADS", "4")
    )
    n_batch = int(
        current_app.config.get("LOCAL_EMBEDDING_N_BATCH", 64)
        or os.getenv("LOCAL_EMBEDDING_N_BATCH", "64")
    )

    return Llama(
        model_path=model_path,
        embedding=True,
        n_ctx=n_ctx,
        n_threads=n_threads,
        n_batch=n_batch,
    )


def _delete_existing_embeddings(document_id):
    return DocumentEmbedding.query.filter_by(document_id=document_id).delete(
        synchronize_session=False
    )


def _embed_document_chunks(document, chunks, embed_fn, metadata=None):
    existing = _delete_existing_embeddings(document.id)
    if existing:
        click.echo(f"Removed existing embeddings: {existing}")

    created = 0
    batch = []

    with click.progressbar(chunks, label="Embedding chunks") as chunk_iter:
        for idx, chunk in enumerate(chunk_iter):
            embedding = embed_fn(chunk)
            batch.append(
                DocumentEmbedding(
                    document_id=document.id,
                    document_type=document.document_type,
                    embedding=embedding,
                    chunk_index=idx,
                    content=chunk,
                    metadata_=metadata if isinstance(metadata, dict) else None,
                )
            )
            created += 1

            if len(batch) >= DEFAULT_BATCH_SIZE:
                db.session.add_all(batch)
                db.session.commit()
                batch = []

    if batch:
        db.session.add_all(batch)
        db.session.commit()

    click.echo(f"Rows produced: {created}")


def _embed_document_from_bytes(
    document,
    data,
    chunk_size,
    chunk_overlap,
    embedder,
    metadata=None,
):
    text_content = _extract_text(document, data)
    if not text_content.strip():
        raise click.ClickException("No text content extracted from document.")

    if embedder["type"] == "openai":
        encoding = _encoding_for_model(embedder["model"])
    else:
        encoding = _encoding_for_model("cl100k_base")
    tokens = encoding.encode(text_content)
    chunks = _chunk_tokens(encoding, tokens, chunk_size, chunk_overlap)

    click.echo(f"Total tokens: {len(tokens)}")
    click.echo(f"Expected chunks: {len(chunks)}")

    if embedder["type"] == "openai":
        client = embedder["client"]
        model = embedder["model"]

        def _openai_embed(chunk):
            response = client.embeddings.create(model=model, input=chunk)
            return response.data[0].embedding

        _embed_document_chunks(document, chunks, _openai_embed, metadata=metadata)
        return

    if embedder["type"] == "local":
        llm = embedder["llm"]

        def _local_embed(chunk):
            response = llm.create_embedding(chunk)
            return response["data"][0]["embedding"]

        _embed_document_chunks(document, chunks, _local_embed, metadata=metadata)
        return

    raise click.ClickException("Unsupported embedder type.")


def _build_sqs_client():
    region = current_app.config.get("SQS_REGION") or os.getenv("SQS_REGION") or None
    endpoint = current_app.config.get("AWS_SQS_ENDPOINT") or os.getenv(
        "AWS_SQS_ENDPOINT"
    )
    return boto3.client("sqs", region_name=region, endpoint_url=endpoint or None)


def _build_s3_client():
    region = current_app.config.get("AWS_REGION") or os.getenv("AWS_REGION") or None
    endpoint = current_app.config.get("AWS_S3_ENDPOINT") or os.getenv("AWS_S3_ENDPOINT")
    access_key = current_app.config.get("AWS_ACCESS_KEY_ID") or os.getenv(
        "AWS_ACCESS_KEY_ID"
    )
    secret_key = current_app.config.get("AWS_SECRET_ACCESS_KEY") or os.getenv(
        "AWS_SECRET_ACCESS_KEY"
    )
    return boto3.client(
        "s3",
        region_name=region,
        endpoint_url=endpoint or None,
        aws_access_key_id=access_key or None,
        aws_secret_access_key=secret_key or None,
    )


def _parse_payload(body):
    try:
        payload = json.loads(body)
    except json.JSONDecodeError as exc:
        raise click.ClickException("SQS message body is not valid JSON.") from exc
    if isinstance(payload, dict) and "Message" in payload and isinstance(
        payload["Message"], str
    ):
        try:
            payload = json.loads(payload["Message"])
        except json.JSONDecodeError:
            return payload
    return payload


def _extract_s3_location(payload):
    if not isinstance(payload, dict):
        return None, None

    bucket = payload.get("bucket") or payload.get("Bucket") or payload.get("s3_bucket")
    key = payload.get("key") or payload.get("Key") or payload.get("s3_key")

    if bucket and key:
        return bucket, key

    records = payload.get("Records")
    if isinstance(records, list) and records:
        record = records[0]
        s3 = record.get("s3", {})
        bucket = s3.get("bucket", {}).get("name")
        key = s3.get("object", {}).get("key")
        if bucket and key:
            return bucket, key

    return None, None


def _unique_document_name(base_name):
    name = base_name
    counter = 1
    while Document.query.filter_by(name=name).first() is not None:
        counter += 1
        name = f"{base_name}-{counter}"
    return name


def _get_or_create_document_from_payload(payload, storage_key, content_type, size_bytes):
    document_id = payload.get("document_id") if isinstance(payload, dict) else None
    if document_id:
        document = db.session.get(Document, document_id)
        if document is None:
            raise click.ClickException("Document not found for payload document_id.")
        return document

    document = Document.query.filter_by(storage_key=storage_key).first()
    if document is not None:
        return document

    filename = payload.get("original_filename") if isinstance(payload, dict) else None
    if not filename:
        filename = os.path.basename(storage_key)
    base_name = payload.get("name") if isinstance(payload, dict) else None
    if not base_name:
        base_name = os.path.splitext(filename)[0] or storage_key

    document = Document(
        name=_unique_document_name(base_name),
        description=payload.get("description") if isinstance(payload, dict) else None,
        document_type=payload.get("document_type") if isinstance(payload, dict) else None,
        original_filename=filename,
        storage_key=storage_key,
        storage_provider="s3",
        content_type=content_type,
        size_bytes=size_bytes,
        embedding_status="pending",
    )
    db.session.add(document)
    db.session.commit()
    return document


def _set_document_status(document, status, enqueue_error=None, embedding_error=None):
    document.embedding_status = status
    if enqueue_error is not None:
        document.enqueue_error = enqueue_error
    if embedding_error is not None:
        document.embedding_error = embedding_error
    db.session.commit()


def register_cli(app):
    @app.cli.group()
    def system():
        """System tasks."""
        pass

    @system.command("greet")
    def greet():
        """Print hello world."""
        click.echo("hello world")

    @system.command("create-admin")
    @click.option("--email", default="admin@example.com", show_default=True)
    @click.option("--password", default="password", show_default=True)
    @click.option("--first-name", default="Test", show_default=True)
    @click.option("--last-name", default="Admin", show_default=True)
    @click.option(
        "--force",
        is_flag=True,
        help="Update existing user if the email already exists.",
    )
    @with_appcontext
    def create_admin(email, password, first_name, last_name, force):
        """Create a default admin user."""
        user = User.query.filter_by(email=email).first()
        document_types = current_app.config.get("DOCUMENT_TYPES") or []
        if user is not None and not force:
            click.echo("Admin user already exists. Use --force to update it.")
            return

        if user is None:
            user = User(
                email=email,
                first_name=first_name,
                last_name=last_name,
                password_hash=build_password_hash(password),
                status="active",
                user_type="admin",
                document_types=document_types,
            )
            db.session.add(user)
            db.session.commit()
            click.echo("Admin user created.")
            return

        user.first_name = first_name
        user.last_name = last_name
        user.password_hash = build_password_hash(password)
        user.status = "active"
        user.user_type = "admin"
        user.document_types = document_types
        db.session.commit()
        click.echo("Admin user updated.")

    @system.command("create-user")
    @click.option("--email", default="user@example.com", show_default=True)
    @click.option("--password", default="password", show_default=True)
    @click.option("--first-name", default="Sample", show_default=True)
    @click.option("--last-name", default="User", show_default=True)
    @click.option(
        "--force",
        is_flag=True,
        help="Update existing user if the email already exists.",
    )
    @with_appcontext
    def create_user(email, password, first_name, last_name, force):
        """Create a default sample user."""
        user = User.query.filter_by(email=email).first()
        document_types = current_app.config.get("DOCUMENT_TYPES") or []
        if user is not None and not force:
            click.echo("User already exists. Use --force to update it.")
            return

        if user is None:
            user = User(
                email=email,
                first_name=first_name,
                last_name=last_name,
                password_hash=build_password_hash(password),
                status="active",
                user_type="user",
                document_types=document_types,
            )
            db.session.add(user)
            db.session.commit()
            click.echo("User created.")
            return

        user.first_name = first_name
        user.last_name = last_name
        user.password_hash = build_password_hash(password)
        user.status = "active"
        user.user_type = "user"
        user.document_types = document_types
        db.session.commit()
        click.echo("User updated.")

    @system.command("create-ops")
    @click.option("--email", default="ops@example.com", show_default=True)
    @click.option("--password", default="password", show_default=True)
    @click.option("--first-name", default="Sample", show_default=True)
    @click.option("--last-name", default="Ops", show_default=True)
    @click.option(
        "--force",
        is_flag=True,
        help="Update existing user if the email already exists.",
    )
    @with_appcontext
    def create_ops(email, password, first_name, last_name, force):
        """Create a default ops user."""
        user = User.query.filter_by(email=email).first()
        document_types = current_app.config.get("DOCUMENT_TYPES") or []
        if user is not None and not force:
            click.echo("Ops user already exists. Use --force to update it.")
            return

        if user is None:
            user = User(
                email=email,
                first_name=first_name,
                last_name=last_name,
                password_hash=build_password_hash(password),
                status="active",
                user_type="ops",
                document_types=document_types,
            )
            db.session.add(user)
            db.session.commit()
            click.echo("Ops user created.")
            return

        user.first_name = first_name
        user.last_name = last_name
        user.password_hash = build_password_hash(password)
        user.status = "active"
        user.user_type = "ops"
        user.document_types = document_types
        db.session.commit()
        click.echo("Ops user updated.")

    @system.command("openai-embed-document")
    @click.option("--document-id", required=True)
    @click.option("--chunk-size", default=DEFAULT_CHUNK_TOKENS, show_default=True, type=int)
    @click.option(
        "--chunk-overlap", default=DEFAULT_CHUNK_OVERLAP, show_default=True, type=int
    )
    @with_appcontext
    def openai_embed_document(document_id, chunk_size, chunk_overlap):
        """Embed a document from storage and save vectors."""
        _validate_chunking_options(chunk_size, chunk_overlap)

        document = db.session.get(Document, document_id)
        if document is None:
            raise click.ClickException("Document not found.")

        click.echo(f"Downloading document {document.id}...")
        data = get_storage().read(document.storage_key)
        client, model = _openai_client_and_model()
        _embed_document_from_bytes(
            document,
            data,
            chunk_size,
            chunk_overlap,
            {"type": "openai", "client": client, "model": model},
        )

    @system.command("local-embed-document")
    @click.option("--document-id", required=True)
    @click.option("--chunk-size", default=DEFAULT_CHUNK_TOKENS, show_default=True, type=int)
    @click.option(
        "--chunk-overlap", default=DEFAULT_CHUNK_OVERLAP, show_default=True, type=int
    )
    @with_appcontext
    def local_embed_document(document_id, chunk_size, chunk_overlap):
        """Embed a document from storage using a local GGUF model."""
        _validate_chunking_options(chunk_size, chunk_overlap)

        document = db.session.get(Document, document_id)
        if document is None:
            raise click.ClickException("Document not found.")

        click.echo(f"Downloading document {document.id}...")
        data = get_storage().read(document.storage_key)
        llm = _local_embedder()
        _embed_document_from_bytes(
            document,
            data,
            chunk_size,
            chunk_overlap,
            {"type": "local", "llm": llm},
        )

    @system.command("process-sqs-embedding")
    @click.option("--queue-url", default=None)
    @click.option("--wait-time", default=DEFAULT_SQS_WAIT_TIME, show_default=True, type=int)
    @click.option("--visibility-timeout", default=120, show_default=True, type=int)
    @click.option("--delete-message/--no-delete-message", default=True, show_default=True)
    @click.option("--max-messages", default=1, show_default=True, type=int)
    @click.option(
        "--embedder",
        type=click.Choice(["auto", "openai", "local"], case_sensitive=False),
        default="auto",
        show_default=True,
    )
    @click.option("--chunk-size", default=DEFAULT_CHUNK_TOKENS, show_default=True, type=int)
    @click.option(
        "--chunk-overlap", default=DEFAULT_CHUNK_OVERLAP, show_default=True, type=int
    )
    @with_appcontext
    def process_sqs_embedding(
        queue_url,
        wait_time,
        visibility_timeout,
        delete_message,
        max_messages,
        embedder,
        chunk_size,
        chunk_overlap,
    ):
        """Poll SQS jobs and embed referenced documents."""
        _validate_chunking_options(chunk_size, chunk_overlap)
        queue_url = queue_url or current_app.config.get("SQS_QUEUE_URL") or os.getenv(
            "SQS_QUEUE_URL"
        )
        if not queue_url:
            raise click.ClickException("SQS_QUEUE_URL is required.")
        if max_messages <= 0 or max_messages > 10:
            raise click.ClickException("--max-messages must be between 1 and 10.")

        sqs = _build_sqs_client()
        click.echo("Polling SQS. Press Ctrl+C to stop.")

        try:
            while True:
                response = sqs.receive_message(
                    QueueUrl=queue_url,
                    MaxNumberOfMessages=max_messages,
                    WaitTimeSeconds=wait_time,
                    VisibilityTimeout=visibility_timeout,
                )
                messages = response.get("Messages") or []
                if not messages:
                    continue

                for message in messages:
                    payload = _parse_payload(message.get("Body", "{}"))
                    try:
                        document = None
                        bucket = None
                        key = None

                        if isinstance(payload, dict) and payload.get("document_id"):
                            document = db.session.get(Document, payload["document_id"])
                            if document is None:
                                click.echo("Skipping message: document not found.")
                                continue
                            key = payload.get("key") or document.storage_key
                            bucket = payload.get("bucket") or current_app.config.get(
                                "AWS_S3_BUCKET"
                            ) or os.getenv("AWS_S3_BUCKET")
                            if not key or not bucket:
                                click.echo("Skipping message: missing bucket/key for document.")
                                continue
                        else:
                            bucket, key = _extract_s3_location(payload)
                            if not bucket or not key:
                                click.echo("Skipping message: payload missing S3 bucket/key.")
                                continue

                        configured_bucket = current_app.config.get("AWS_S3_BUCKET") or os.getenv(
                            "AWS_S3_BUCKET"
                        )
                        if configured_bucket and bucket != configured_bucket:
                            click.echo(
                                "Skipping message: payload bucket does not match AWS_S3_BUCKET."
                            )
                            continue

                        key = unquote_plus(key)
                        click.echo(f"Downloading s3://{bucket}/{key} ...")
                        s3 = _build_s3_client()
                        obj = s3.get_object(Bucket=bucket, Key=key)
                        data = obj["Body"].read()
                        content_type = obj.get("ContentType")
                        size_bytes = obj.get("ContentLength")

                        if document is None:
                            document = _get_or_create_document_from_payload(
                                payload, key, content_type, size_bytes
                            )

                        _set_document_status(document, "processing", embedding_error=None)

                        if embedder.lower() == "openai":
                            client, model = _openai_client_and_model()
                            embedder_config = {"type": "openai", "client": client, "model": model}
                        elif embedder.lower() == "local":
                            embedder_config = {"type": "local", "llm": _local_embedder()}
                        else:
                            use_openai = str(current_app.config.get("USE_OPENAI", "true")).lower()
                            if use_openai in {"1", "true", "yes", "y"}:
                                client, model = _openai_client_and_model()
                                embedder_config = {"type": "openai", "client": client, "model": model}
                            else:
                                embedder_config = {"type": "local", "llm": _local_embedder()}

                        _embed_document_from_bytes(
                            document,
                            data,
                            chunk_size,
                            chunk_overlap,
                            embedder_config,
                            metadata=payload.get("metadata") if isinstance(payload, dict) else None,
                        )

                        _set_document_status(
                            document,
                            "embedded",
                            enqueue_error=None,
                            embedding_error=None,
                        )

                        if delete_message:
                            sqs.delete_message(
                                QueueUrl=queue_url, ReceiptHandle=message["ReceiptHandle"]
                            )
                            click.echo("Deleted SQS message.")
                    except Exception as exc:  # noqa: BLE001 - keep worker running
                        if "document" in locals() and document is not None:
                            _set_document_status(document, "failed", embedding_error=str(exc))
                        click.echo(f"Embedding failed: {exc}")
        except KeyboardInterrupt:
            click.echo("Shutting down SQS poller.")


def _quote_identifier(name):
    return f"\"{name.replace('\"', '\"\"')}\""


def _ensure_sqlite_db(database_path):
    if not database_path or database_path == ":memory:":
        click.echo("SQLite in-memory database does not need creation.")
        return
    path = Path(database_path)
    if not path.is_absolute():
        path = Path.cwd() / path
    path.parent.mkdir(parents=True, exist_ok=True)
    engine = create_engine(f"sqlite:///{path}")
    engine.connect().close()
    click.echo(f"SQLite database ready at {path}")


@db_cli.command("create")
@with_appcontext
def create_db():
    """Create the configured database if it does not exist."""
    database_uri = current_app.config.get("SQLALCHEMY_DATABASE_URI")
    if not database_uri:
        raise click.ClickException("SQLALCHEMY_DATABASE_URI is not configured.")

    url = make_url(database_uri)
    if url.get_backend_name() == "sqlite":
        _ensure_sqlite_db(url.database)
        return

    db_name = url.database
    if not db_name:
        raise click.ClickException("Database name is missing from SQLALCHEMY_DATABASE_URI.")

    try:
        admin_url = url.set(database="postgres")
    except AttributeError:
        admin_url = url._replace(database="postgres")

    engine = create_engine(admin_url, isolation_level="AUTOCOMMIT")
    with engine.connect() as connection:
        exists = connection.execute(
            text("SELECT 1 FROM pg_database WHERE datname = :name"),
            {"name": db_name},
        ).scalar()
        if exists:
            click.echo(f"Database already exists: {db_name}")
            return
        connection.execute(text(f"CREATE DATABASE {_quote_identifier(db_name)}"))
        click.echo(f"Database created: {db_name}")


def _extract_text(document, data):
    filename = (document.original_filename or "").lower()
    content_type = (document.content_type or "").lower()
    is_pdf = (
        content_type == "application/pdf"
        or filename.endswith(".pdf")
        or data[:4] == b"%PDF"
    )
    is_excel = (
        content_type
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or filename.endswith(".xlsx")
    )
    is_powerpoint = (
        content_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or filename.endswith(".pptx")
    )

    if is_pdf:
        reader = PdfReader(io.BytesIO(data))
        parts = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(parts)

    if is_excel:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for sheet in workbook.worksheets:
            parts.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(value) for value in row if value is not None]
                if row_values:
                    parts.append("\t".join(row_values))
        return "\n".join(parts)

    if is_powerpoint:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        parts = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts.append(f"# Slide {slide_index}")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text:
                    parts.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        row_values = [
                            cell.text for cell in row.cells if cell.text
                        ]
                        if row_values:
                            parts.append("\t".join(row_values))
        return "\n".join(parts)

    return data.decode("utf-8", errors="ignore")


def _encoding_for_model(model):
    try:
        return tiktoken.encoding_for_model(model)
    except KeyError:
        return tiktoken.get_encoding("cl100k_base")


def _chunk_tokens(encoding, tokens, chunk_size, chunk_overlap):
    if not tokens:
        return []
    step = max(chunk_size - chunk_overlap, 1)
    return [
        encoding.decode(tokens[i : i + chunk_size])
        for i in range(0, len(tokens), step)
    ]
